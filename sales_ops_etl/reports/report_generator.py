"""
report_generator.py
-------------------
Generates two stakeholder-ready outputs:
  1. Excel workbook with multiple KPI sheets + formatted tables
  2. PowerPoint executive summary with key metrics and charts
"""

import os
import sys
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
from openpyxl import load_workbook
from openpyxl.styles import (PatternFill, Font, Alignment, Border, Side,
                              numbers as xl_numbers)
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, LineChart, Reference
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import warnings
warnings.filterwarnings('ignore')

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'analysis'))
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'etl'))
import kpi_calculator
import load

EXPORT_DIR = os.path.join(os.path.dirname(__file__), '..', 'data', 'exports')
os.makedirs(EXPORT_DIR, exist_ok=True)

# Brand colours
BLUE   = '2563EB'
GREEN  = '10B981'
ORANGE = 'F59E0B'
RED    = 'EF4444'
GREY   = 'F1F5F9'
DARK   = '1E293B'


# ── Helpers ───────────────────────────────────────────────────

def _header_fill(ws, row: int, ncols: int, hex_color: str = BLUE):
    fill = PatternFill(fill_type='solid', fgColor=hex_color)
    font = Font(color='FFFFFF', bold=True, size=11)
    for col in range(1, ncols + 1):
        cell = ws.cell(row=row, column=col)
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(horizontal='center', vertical='center')


def _auto_width(ws):
    for col in ws.columns:
        max_len = max((len(str(c.value)) if c.value else 0 for c in col), default=10)
        ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 35)


def _stripe_rows(ws, start_row: int, end_row: int, ncols: int):
    light = PatternFill(fill_type='solid', fgColor=GREY)
    for r in range(start_row, end_row + 1):
        if r % 2 == 0:
            for c in range(1, ncols + 1):
                ws.cell(row=r, column=c).fill = light


# ── Excel Report ──────────────────────────────────────────────

def build_excel_report(kpis: dict) -> str:
    path = os.path.join(EXPORT_DIR, 'Sales_Ops_KPI_Report.xlsx')

    with pd.ExcelWriter(path, engine='openpyxl') as writer:

        # Sheet 1: Executive Summary
        mom = kpis['mom_variance']
        total_rev   = load.query("SELECT SUM(revenue) AS r FROM fact_revenue")['r'].iloc[0]
        avg_wr      = kpis['win_rate']['win_rate'].mean()
        avg_attain  = kpis['quota_attainment']['quota_attainment'].mean()
        avg_acc     = kpis['forecast_accuracy']['forecast_accuracy'].mean()

        summary = pd.DataFrame({
            'KPI':   ['Total Revenue', 'Avg Win Rate', 'Avg Quota Attainment',
                      'Avg Forecast Accuracy', 'Total Closed Deals'],
            'Value': [f"${total_rev:,.0f}", f"{avg_wr:.1%}", f"{avg_attain:.1%}",
                      f"{avg_acc:.1%}",
                      str(kpis['win_rate']['won'].sum())]
        })
        summary.to_excel(writer, sheet_name='Executive Summary', index=False)

        # Sheet 2: Win Rate
        kpis['win_rate'].to_excel(writer, sheet_name='Win Rate', index=False)

        # Sheet 3: Pipeline Velocity
        kpis['pipeline_velocity'].to_excel(writer, sheet_name='Pipeline Velocity', index=False)

        # Sheet 4: Quota Attainment
        kpis['quota_attainment'].drop(columns=['attainment_band'], errors='ignore') \
                                .to_excel(writer, sheet_name='Quota Attainment', index=False)

        # Sheet 5: Forecast Accuracy
        kpis['forecast_accuracy'].to_excel(writer, sheet_name='Forecast Accuracy', index=False)

        # Sheet 6: MoM Variance
        kpis['mom_variance'].to_excel(writer, sheet_name='MoM Revenue', index=False)

        # Sheet 7: Territory
        kpis['territory_perf'].to_excel(writer, sheet_name='Territory Performance', index=False)

    # Post-process: style with openpyxl
    wb = load_workbook(path)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        ncols = ws.max_column
        _header_fill(ws, 1, ncols)
        _stripe_rows(ws, 2, ws.max_row, ncols)
        _auto_width(ws)
    wb.save(path)
    return path


# ── Charts for PowerPoint ─────────────────────────────────────

def _make_revenue_chart(kpis: dict) -> str:
    mom = kpis['mom_variance'].tail(12)
    fig, ax = plt.subplots(figsize=(9, 4))
    bars = ax.bar(mom['month_year'], mom['total_revenue'] / 1e6,
                  color=['#10B981' if v >= 0 else '#EF4444' for v in mom['variance']],
                  edgecolor='white', linewidth=0.5)
    ax.set_title('Monthly Revenue ($M)', fontsize=13, fontweight='bold', color='#1E293B')
    ax.set_xlabel('')
    ax.yaxis.set_major_formatter(mtick.FormatStrFormatter('$%.1fM'))
    ax.tick_params(axis='x', rotation=45, labelsize=8)
    ax.spines[['top', 'right']].set_visible(False)
    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    p = os.path.join(EXPORT_DIR, '_chart_revenue.png')
    fig.savefig(p, dpi=150, bbox_inches='tight')
    plt.close()
    return p


def _make_winrate_chart(kpis: dict) -> str:
    wr = kpis['win_rate'].groupby('month_year')['win_rate'].mean().reset_index()
    wr = wr.tail(12)
    fig, ax = plt.subplots(figsize=(9, 4))
    ax.plot(wr['month_year'], wr['win_rate'] * 100, color='#2563EB',
            linewidth=2.5, marker='o', markersize=5)
    ax.axhline(wr['win_rate'].mean() * 100, color='#F59E0B',
               linestyle='--', linewidth=1.5, label=f"Avg: {wr['win_rate'].mean():.1%}")
    ax.set_title('Monthly Win Rate (%)', fontsize=13, fontweight='bold', color='#1E293B')
    ax.yaxis.set_major_formatter(mtick.PercentFormatter())
    ax.tick_params(axis='x', rotation=45, labelsize=8)
    ax.spines[['top', 'right']].set_visible(False)
    ax.legend(fontsize=9)
    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    p = os.path.join(EXPORT_DIR, '_chart_winrate.png')
    fig.savefig(p, dpi=150, bbox_inches='tight')
    plt.close()
    return p


def _make_territory_chart(kpis: dict) -> str:
    t = kpis['territory_perf'].sort_values('total_revenue', ascending=True).tail(8)
    fig, ax = plt.subplots(figsize=(9, 4))
    bars = ax.barh(t['territory'], t['total_revenue'] / 1e6, color='#2563EB', alpha=0.85)
    ax.set_title('Revenue by Territory ($M)', fontsize=13, fontweight='bold', color='#1E293B')
    ax.xaxis.set_major_formatter(mtick.FormatStrFormatter('$%.1fM'))
    ax.spines[['top', 'right']].set_visible(False)
    ax.grid(axis='x', alpha=0.3)
    plt.tight_layout()
    p = os.path.join(EXPORT_DIR, '_chart_territory.png')
    fig.savefig(p, dpi=150, bbox_inches='tight')
    plt.close()
    return p


# ── PowerPoint Report ─────────────────────────────────────────

def build_pptx_report(kpis: dict, chart_paths: dict) -> str:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    def add_textbox(slide, text, left, top, width, height,
                    size=14, bold=False, color='1E293B', align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)

    def add_rect(slide, left, top, width, height, fill_hex='2563EB'):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
        shape.line.fill.background()
        return shape

    # ── Slide 1: Title ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, 'F8FAFC')
    add_rect(slide, 0, 0, 13.33, 1.5, BLUE)
    add_textbox(slide, 'Sales Operations', 0.5, 0.15, 12, 0.7,
                size=32, bold=True, color='FFFFFF')
    add_textbox(slide, 'KPI Reporting & Analytics Dashboard', 0.5, 0.85, 12, 0.6,
                size=18, color='BFDBFE')
    add_textbox(slide, 'Automated ETL Pipeline  |  Statistical Analysis  |  Executive Reporting',
                0.5, 1.8, 12, 0.5, size=13, color='64748B')

    total_rev  = load.query("SELECT SUM(revenue) AS r FROM fact_revenue")['r'].iloc[0]
    avg_wr     = kpis['win_rate']['win_rate'].mean()
    avg_att    = kpis['quota_attainment']['quota_attainment'].mean()

    metrics = [
        (f"${total_rev/1e6:.1f}M", "Total Revenue"),
        (f"{avg_wr:.1%}",          "Avg Win Rate"),
        (f"{avg_att:.1%}",         "Avg Quota Attainment"),
        (f"{len(kpis['territory_perf'])}",   "Territories Tracked"),
    ]
    for i, (val, lbl) in enumerate(metrics):
        x = 1.0 + i * 3.0
        add_rect(slide, x, 2.6, 2.6, 1.5, 'EFF6FF')
        add_textbox(slide, val, x, 2.75, 2.6, 0.7, size=28, bold=True,
                    color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, x, 3.4, 2.6, 0.4, size=11,
                    color='64748B', align=PP_ALIGN.CENTER)

    add_textbox(slide, 'Period: Jan 2023 – Dec 2024  |  Data Source: CRM + ERP + Marketing DB',
                0.5, 6.8, 12, 0.4, size=9, color='94A3B8', align=PP_ALIGN.CENTER)

    # ── Slide 2: Revenue Trend ────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 0.7, BLUE)
    add_textbox(slide, 'Monthly Revenue Performance', 0.3, 0.1, 10, 0.5,
                size=20, bold=True, color='FFFFFF')
    slide.shapes.add_picture(chart_paths['revenue'], Inches(0.3), Inches(0.9),
                             Inches(9.0), Inches(4.0))
    mom = kpis['mom_variance']
    avg_growth = mom['variance_pct'].mean()
    add_rect(slide, 9.8, 0.9, 3.2, 1.2, 'EFF6FF')
    add_textbox(slide, f"{avg_growth:+.1f}%", 9.8, 1.0, 3.2, 0.6,
                size=26, bold=True, color=GREEN if avg_growth >= 0 else RED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 'Avg MoM Growth', 9.8, 1.55, 3.2, 0.4, size=11,
                color='64748B', align=PP_ALIGN.CENTER)

    # ── Slide 3: Win Rate ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 0.7, '0F172A')
    add_textbox(slide, 'Win Rate Analysis', 0.3, 0.1, 10, 0.5,
                size=20, bold=True, color='FFFFFF')
    slide.shapes.add_picture(chart_paths['winrate'], Inches(0.3), Inches(0.9),
                             Inches(9.0), Inches(4.0))
    top_rep = kpis['win_rate'].groupby('rep_name')['win_rate'].mean().idxmax()
    top_wr  = kpis['win_rate'].groupby('rep_name')['win_rate'].mean().max()
    add_rect(slide, 9.8, 0.9, 3.2, 1.8, 'F0FDF4')
    add_textbox(slide, 'Top Performer', 9.8, 1.0, 3.2, 0.4, size=10,
                color='64748B', align=PP_ALIGN.CENTER)
    add_textbox(slide, top_rep.split()[-1], 9.8, 1.35, 3.2, 0.5,
                size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, f"{top_wr:.1%} win rate", 9.8, 1.85, 3.2, 0.4,
                size=13, color=GREEN, align=PP_ALIGN.CENTER)

    # ── Slide 4: Territory Performance ───────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 0.7, ORANGE)
    add_textbox(slide, 'Territory Performance', 0.3, 0.1, 10, 0.5,
                size=20, bold=True, color='FFFFFF')
    slide.shapes.add_picture(chart_paths['territory'], Inches(0.3), Inches(0.9),
                             Inches(9.0), Inches(4.0))

    # ── Slide 5: Key Findings ─────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, 'F8FAFC')
    add_rect(slide, 0, 0, 13.33, 0.8, DARK)
    add_textbox(slide, 'Key Findings & Recommendations', 0.3, 0.1, 12, 0.65,
                size=22, bold=True, color='FFFFFF')

    findings = [
        (GREEN,  "Win Rate",         f"Average {avg_wr:.1%} — exceeds industry benchmark of 25–30%."),
        (BLUE,   "Quota Attainment", f"Team avg {kpis['quota_attainment']['quota_attainment'].mean():.1%} — focus coaching on reps below 80%."),
        (ORANGE, "Forecast Accuracy",f"Avg MAPE {kpis['forecast_accuracy']['mape'].mean():.1f}% — refine inputs for outlier months."),
        (RED,    "Sales Cycle",      f"Avg {kpis['territory_perf']['avg_cycle_days'].mean():.0f} days — identify bottlenecks in Proposal stage."),
    ]
    for i, (color, title, text) in enumerate(findings):
        y = 1.1 + i * 1.4
        add_rect(slide, 0.4, y, 0.08, 0.9, color)
        add_textbox(slide, title, 0.65, y,        11, 0.42, size=14, bold=True, color=DARK)
        add_textbox(slide, text,  0.65, y + 0.42, 11, 0.5,  size=12, color='475569')

    path = os.path.join(EXPORT_DIR, 'Sales_Ops_Executive_Report.pptx')
    prs.save(path)
    return path


# ── Main ──────────────────────────────────────────────────────

def run():
    print("── Generating Reports ───────────────────────────────────")

    kpis = kpi_calculator.run()
    print()

    chart_paths = {
        'revenue':   _make_revenue_chart(kpis),
        'winrate':   _make_winrate_chart(kpis),
        'territory': _make_territory_chart(kpis),
    }
    print("  Charts generated")

    excel_path = build_excel_report(kpis)
    print(f"  Excel report:      {excel_path.split('/')[-1]}")

    pptx_path = build_pptx_report(kpis, chart_paths)
    print(f"  PowerPoint report: {pptx_path.split('/')[-1]}")

    return excel_path, pptx_path


if __name__ == '__main__':
    run()
